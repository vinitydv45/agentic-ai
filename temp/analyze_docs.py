#!/usr/bin/env python3
"""Analyze PowerPoint and Word documents to understand their structure."""

from pptx import Presentation
from docx import Document
import json

def analyze_pptx(file_path):
    """Analyze PowerPoint file structure."""
    print(f"\n{'='*80}")
    print(f"ANALYZING: {file_path}")
    print(f"{'='*80}\n")

    prs = Presentation(file_path)

    print(f"Total Slides: {len(prs.slides)}\n")

    for idx, slide in enumerate(prs.slides):
        print(f"\n{'─'*80}")
        print(f"SLIDE {idx + 1}")
        print(f"{'─'*80}")
        print(f"Layout: {slide.slide_layout.name}")
        print(f"Shapes: {len(slide.shapes)}\n")

        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                print(f"  Shape {shape_idx + 1} ({shape.shape_type}): {shape.text[:100]}...")
            elif hasattr(shape, "table"):
                print(f"  Shape {shape_idx + 1}: TABLE ({shape.table.rows.__len__()} rows x {shape.table.columns.__len__()} cols)")

    return prs

def analyze_docx(file_path):
    """Analyze Word document structure."""
    print(f"\n{'='*80}")
    print(f"ANALYZING: {file_path}")
    print(f"{'='*80}\n")

    doc = Document(file_path)

    print(f"Total Paragraphs: {len(doc.paragraphs)}")
    print(f"Total Tables: {len(doc.tables)}\n")

    print(f"\n{'─'*80}")
    print("CONTENT STRUCTURE")
    print(f"{'─'*80}\n")

    for idx, para in enumerate(doc.paragraphs):
        if para.text.strip():
            style = para.style.name
            text = para.text[:150]
            print(f"Para {idx + 1} [{style}]: {text}...")

    if doc.tables:
        print(f"\n{'─'*80}")
        print("TABLES")
        print(f"{'─'*80}\n")

        for table_idx, table in enumerate(doc.tables):
            print(f"\nTable {table_idx + 1}: {len(table.rows)} rows x {len(table.columns)} cols")

    return doc

if __name__ == "__main__":
    # Analyze PowerPoint
    pptx_file = "Monthly-Connect_template_new.pptx"
    prs = analyze_pptx(pptx_file)

    # Analyze Word document
    docx_file = "Samsung_PRISM_Student_Project_Handbook.docx"
    doc = analyze_docx(docx_file)

    print(f"\n{'='*80}")
    print("ANALYSIS COMPLETE")
    print(f"{'='*80}\n")
